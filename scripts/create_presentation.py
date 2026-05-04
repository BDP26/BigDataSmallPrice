"""Generate BigDataSmallPrice.pptx with modern dark theme."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colours ────────────────────────────────────────────────────────────────────
BG          = RGBColor(0x0a, 0x0f, 0x1e)   # deep navy
SURFACE     = RGBColor(0x11, 0x18, 0x27)
CARD        = RGBColor(0x1a, 0x22, 0x36)
BLUE        = RGBColor(0x3b, 0x82, 0xf6)
CYAN        = RGBColor(0x06, 0xb6, 0xd4)
GREEN       = RGBColor(0x10, 0xb9, 0x81)
YELLOW      = RGBColor(0xf5, 0x9e, 0x0b)
RED         = RGBColor(0xef, 0x44, 0x44)
WHITE       = RGBColor(0xff, 0xff, 0xff)
MUTED       = RGBColor(0x94, 0xa3, 0xb8)
BORDER      = RGBColor(0x2a, 0x3a, 0x5a)
LIGHT_BLUE  = RGBColor(0x60, 0xa5, 0xfa)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # rectangle = 1
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Calibri"
    return txb


def add_slide(title_text, label_text=None):
    slide = prs.slides.add_slide(blank_layout)
    # Dark background
    bg_shape = add_rect(slide, 0, 0, W, H, fill_color=BG)
    bg_shape.zorder = 0

    # Top accent line
    add_rect(slide, 0, 0, W, Pt(3), fill_color=BLUE, line_color=None, line_width=Pt(0))

    # Slide label (top-right, subtle)
    return slide


def set_bg(slide, color=BG):
    """Fill slide background via XML."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def label(slide, text, left=Inches(0.6), top=Inches(0.45)):
    add_text(slide, text.upper(), left, top, Inches(8), Inches(0.3),
             font_size=Pt(9), bold=True, color=CYAN)


def heading(slide, text, left=Inches(0.6), top=Inches(0.8), width=Inches(12), size=Pt(32), color=WHITE):
    add_text(slide, text, left, top, width, Inches(1.2),
             font_size=size, bold=True, color=color)


def body(slide, text, left, top, width, height, size=Pt(13), color=MUTED, align=PP_ALIGN.LEFT):
    add_text(slide, text, left, top, width, height,
             font_size=size, color=color, align=align)


def card_box(slide, left, top, width, height, border_color=BORDER):
    r = add_rect(slide, left, top, width, height, fill_color=CARD, line_color=border_color, line_width=Pt(0.75))
    return r


def bullet_text(slide, items, left, top, width, size=Pt(12.5), color=MUTED, bullet_color=CYAN):
    """Add a list of bullet strings as separate text boxes."""
    row_h = Inches(0.38)
    for i, item in enumerate(items):
        y = top + i * row_h
        # bullet symbol
        add_text(slide, "▸", left, y, Inches(0.25), row_h,
                 font_size=size, color=bullet_color)
        add_text(slide, item, left + Inches(0.28), y, width - Inches(0.28), row_h,
                 font_size=size, color=color)


def metric_box(slide, left, top, width, height, value, unit, label_text, color=BLUE):
    card_box(slide, left, top, width, height, border_color=color)
    add_text(slide, value, left + Inches(0.25), top + Inches(0.2), width - Inches(0.3), Inches(0.7),
             font_size=Pt(36), bold=True, color=color)
    add_text(slide, unit,  left + Inches(0.25), top + Inches(0.82), width - Inches(0.3), Inches(0.3),
             font_size=Pt(11), color=MUTED)
    add_text(slide, label_text, left + Inches(0.25), top + Inches(1.1), width - Inches(0.3), Inches(0.3),
             font_size=Pt(10), color=MUTED)


def bar_row(slide, name, pct, value_text, top, color, left=Inches(0.6), track_w=Inches(5.5)):
    name_w = Inches(1.1)
    val_w  = Inches(1.0)
    gap    = Inches(0.1)
    h      = Inches(0.18)
    row_h  = Inches(0.42)
    bar_y  = top + (row_h - h) / 2

    add_text(slide, name, left, top, name_w, row_h,
             font_size=Pt(11), color=MUTED, align=PP_ALIGN.RIGHT)

    # track
    track_x = left + name_w + gap
    add_rect(slide, track_x, bar_y, track_w, h,
             fill_color=RGBColor(0x1e, 0x2d, 0x45), line_color=None, line_width=Pt(0))

    # fill
    fill_w = int(track_w * pct)
    if fill_w > 0:
        add_rect(slide, track_x, bar_y, fill_w, h,
                 fill_color=color, line_color=None, line_width=Pt(0))

    # value
    add_text(slide, value_text, track_x + track_w + gap, top, val_w, row_h,
             font_size=Pt(11), bold=True, color=color)


def flow_step(slide, left, top, w, h, icon, step_label, step_sub, border_color=BORDER):
    card_box(slide, left, top, w, h, border_color=border_color)
    add_text(slide, icon,       left, top + Inches(0.1),  w, Inches(0.35),
             font_size=Pt(20), align=PP_ALIGN.CENTER)
    add_text(slide, step_label, left, top + Inches(0.45), w, Inches(0.28),
             font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, step_sub,   left, top + Inches(0.72), w, Inches(0.22),
             font_size=Pt(8.5), color=MUTED, align=PP_ALIGN.CENTER)


def arrow(slide, left, top):
    add_text(slide, "→", left, top, Inches(0.3), Inches(0.3),
             font_size=Pt(14), color=BORDER, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

# Top accent bar
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
# Bottom accent bar
add_rect(slide, 0, H - Pt(4), W, Pt(4), fill_color=CYAN, line_color=None, line_width=Pt(0))

# Decorative right panel
add_rect(slide, Inches(8.8), 0, Inches(4.53), H,
         fill_color=RGBColor(0x0d, 0x15, 0x26), line_color=None, line_width=Pt(0))

label(slide, "ZHAW · DS23t PM4 · Big Data")

# Big title
add_text(slide, "BigData", Inches(0.6), Inches(1.1), Inches(7.5), Inches(1.4),
         font_size=Pt(72), bold=True, color=BLUE)
add_text(slide, "SmallPrice", Inches(0.6), Inches(2.3), Inches(7.5), Inches(1.4),
         font_size=Pt(72), bold=True, color=CYAN)

add_text(slide, "Dynamische Strompreisprognose für Winterthur",
         Inches(0.6), Inches(3.7), Inches(7.8), Inches(0.5),
         font_size=Pt(17), color=WHITE)
add_text(slide, "mittels XGBoost, ENTSO-E & Echtzeit-Wetterdaten",
         Inches(0.6), Inches(4.1), Inches(7.8), Inches(0.4),
         font_size=Pt(14), color=MUTED, italic=True)

# Team box on right panel
add_text(slide, "Team", Inches(9.1), Inches(1.5), Inches(3.8), Inches(0.4),
         font_size=Pt(10), bold=True, color=CYAN)
for i, (name, col) in enumerate([
    ("Ryan Bachmann",      BLUE),
    ("Miguel Dinis Silva", CYAN),
    ("Gian Ruchti",        GREEN),
]):
    y = Inches(1.95) + i * Inches(0.5)
    add_rect(slide, Inches(9.1), y, Inches(3.8), Inches(0.38),
             fill_color=CARD, line_color=col, line_width=Pt(0.5))
    add_text(slide, name, Inches(9.25), y + Inches(0.05), Inches(3.5), Inches(0.3),
             font_size=Pt(13), bold=True, color=WHITE)

add_text(slide, "Tech Stack", Inches(9.1), Inches(3.4), Inches(3.8), Inches(0.4),
         font_size=Pt(10), bold=True, color=CYAN)
techs = ["Apache Airflow", "TimescaleDB", "XGBoost", "FastAPI", "Docker"]
for i, t in enumerate(techs):
    y = Inches(3.85) + i * Inches(0.42)
    add_text(slide, "▸  " + t, Inches(9.2), y, Inches(3.6), Inches(0.35),
             font_size=Pt(12), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – MOTIVATION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
label(slide, "Motivation")
heading(slide, "Die Schweizer Energiewende schafft neue Anforderungen")

cards = [
    ("⚡", "Ab 2026 gesetzlich verpflichtend",
     "Netzbetreiber müssen Endkunden dynamische Stromtarife anbieten – Preise reagieren stündlich auf Markt und Einspeisung.",
     BLUE),
    ("🌞", "Volatile erneuerbare Energie",
     "PV und Wind führen zu starken Preisschwankungen. Negative Preise werden häufiger.",
     YELLOW),
    ("🏠", "Enormes Sparpotenzial",
     "Haushalte können durch zeitliches Verschieben von Verbrauch signifikant Kosten sparen.",
     GREEN),
    ("📊", "Keine zugängliche Prognose",
     "Es fehlt ein datengetriebenes Prognose­system für dynamische Strompreise für Endkunden.",
     CYAN),
]

col_w  = Inches(2.95)
col_h  = Inches(2.8)
top    = Inches(2.0)
gap    = Inches(0.22)

for i, (icon, title, desc, color) in enumerate(cards):
    left = Inches(0.5) + i * (col_w + gap)
    card_box(slide, left, top, col_w, col_h, border_color=color)
    add_text(slide, icon,  left + Inches(0.2), top + Inches(0.15), col_w, Inches(0.45),
             font_size=Pt(26))
    add_text(slide, title, left + Inches(0.2), top + Inches(0.62), col_w - Inches(0.3), Inches(0.55),
             font_size=Pt(13), bold=True, color=WHITE)
    add_text(slide, desc,  left + Inches(0.2), top + Inches(1.2), col_w - Inches(0.3), Inches(1.4),
             font_size=Pt(11), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – FRAGESTELLUNG
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
label(slide, "Fragestellung & Ziele")
heading(slide, "Was wollten wir erreichen?")

# Quote box
add_rect(slide, Inches(0.6), Inches(1.8), Pt(4), Inches(1.5),
         fill_color=CYAN, line_color=None, line_width=Pt(0))
add_rect(slide, Inches(0.72), Inches(1.75), Inches(12.2), Inches(1.6),
         fill_color=RGBColor(0x06, 0x2a, 0x35), line_color=CYAN, line_width=Pt(0.75))
add_text(slide,
         "Wie genau lassen sich kurzfristige Strompreise für den Standort Winterthur durch die Kombination "
         "von Day-Ahead-Preisen und meteorologischen Echtzeitdaten prognostizieren?",
         Inches(1.0), Inches(1.85), Inches(11.8), Inches(1.4),
         font_size=Pt(14.5), color=WHITE, italic=True)

# Two columns of bullets
left_items = [
    "Automatisierte Datenpipeline für ENTSO-E & Wetterdaten",
    "Model A: Netzlast-Prognose (kWh) → Netzpreiskomponente",
    "Model B: EPEX Day-Ahead Preisprognose (€/MWh)",
    "Vollständige Tarifsberechnung in Rp./kWh",
]
right_items = [
    "Quantitative Evaluation mit echten Marktdaten",
    "REST API mit JWT-Authentifizierung",
    "Interaktives Dashboard für Endkunden & Admins",
    "Handlungsempfehlungen via Ampelsystem",
]

bullet_text(slide, left_items,  Inches(0.6), Inches(3.65), Inches(5.9), size=Pt(13))
bullet_text(slide, right_items, Inches(6.8), Inches(3.65), Inches(5.9), size=Pt(13))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – ARCHITEKTUR
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
label(slide, "System Architektur")
heading(slide, "End-to-End Big Data Pipeline")

steps = [
    ("🌐", "Datenquellen", "6 APIs",           BLUE),
    ("🔄", "ETL / Airflow", "tägl. 06:00 UTC", CYAN),
    ("🗄️", "TimescaleDB",  "Mio. Datenpunkte", GREEN),
    ("⚙️", "Feature Eng.", "Parquet Export",   YELLOW),
    ("🤖", "XGBoost",      "Model A + B",       BLUE),
    ("🚀", "FastAPI",       "REST + JWT",        CYAN),
    ("📱", "Dashboard",    "User + Admin",       GREEN),
]

step_w  = Inches(1.6)
step_h  = Inches(1.15)
top_f   = Inches(2.0)
total_w = len(steps) * step_w + (len(steps) - 1) * Inches(0.22)
start_x = (W - total_w) / 2

for i, (icon, sl, sub, col) in enumerate(steps):
    lx = start_x + i * (step_w + Inches(0.22))
    flow_step(slide, lx, top_f, step_w, step_h, icon, sl, sub, border_color=col)
    if i < len(steps) - 1:
        ax = lx + step_w + Inches(0.03)
        add_text(slide, "→", ax, top_f + Inches(0.38), Inches(0.2), Inches(0.35),
                 font_size=Pt(14), color=BORDER, align=PP_ALIGN.CENTER)

# Info cards below
card_box(slide, Inches(0.6),  Inches(3.6), Inches(5.8), Inches(1.5), border_color=BLUE)
add_text(slide, "🐳  Docker Compose Stack", Inches(0.8), Inches(3.75), Inches(5.4), Inches(0.4),
         font_size=Pt(13), bold=True, color=CYAN)
add_text(slide, "Airflow (Scheduler + Webserver), TimescaleDB, FastAPI, Frontend – alles containerisiert und reproduzierbar.",
         Inches(0.8), Inches(4.15), Inches(5.4), Inches(0.85),
         font_size=Pt(11.5), color=MUTED)

card_box(slide, Inches(6.8),  Inches(3.6), Inches(5.9), Inches(1.5), border_color=GREEN)
add_text(slide, "📦  Parquet / Chronologischer Split", Inches(7.0), Inches(3.75), Inches(5.5), Inches(0.4),
         font_size=Pt(13), bold=True, color=GREEN)
add_text(slide, "Train / Val / Test chronologisch getrennt. Modelle täglich neu trainiert und mit Zeitstempel versioniert.",
         Inches(7.0), Inches(4.15), Inches(5.5), Inches(0.85),
         font_size=Pt(11.5), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – DATENQUELLEN
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
label(slide, "Daten")
heading(slide, "Sechs heterogene Datenquellen")

sources = [
    ("⚡", "ENTSO-E",              "Day-Ahead-Preise, Netzlast, Erzeugung, grenzüberschreitende Flüsse · <400 API-Calls/Tag", BLUE),
    ("🌤️", "Open-Meteo",           "Temperatur, Globalstrahlung, Bewölkung, Wind, Niederschlag · 10 000 Calls/Tag", CYAN),
    ("🏙️", "Stadtwerk Winterthur", "Reale Netzlast-Zeitreihe & PV-Einspeisung für Winterthur (Grundlage Model A)", GREEN),
    ("🏢", "EKZ / CKW / Groupe E", "Dynamische Stromtarife der Schweizer Netzbetreiber als Evaluationsbasis", YELLOW),
    ("💧", "BAFU Hydrologie",       "Wasserpegel von Flüssen & Seen → Indikator für Wasserkraft-Einspeisung", CYAN),
    ("📊", "Gesamtvolumen",         "Mehrere Millionen Datenpunkte in TimescaleDB mit automatischer Komprimierung", BLUE),
]

col_w = Inches(4.0)
col_h = Inches(1.45)
gap_x = Inches(0.23)
gap_y = Inches(0.2)
top_s = Inches(2.0)

for i, (icon, name, desc, color) in enumerate(sources):
    row  = i // 3
    col  = i  % 3
    lx   = Inches(0.5) + col * (col_w + gap_x)
    ly   = top_s + row * (col_h + gap_y)
    card_box(slide, lx, ly, col_w, col_h, border_color=color)
    add_text(slide, icon + "  " + name,
             lx + Inches(0.2), ly + Inches(0.15), col_w - Inches(0.25), Inches(0.38),
             font_size=Pt(13), bold=True, color=WHITE)
    add_text(slide, desc,
             lx + Inches(0.2), ly + Inches(0.55), col_w - Inches(0.25), Inches(0.78),
             font_size=Pt(10.5), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – METHODE / FEATURE ENGINEERING
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
label(slide, "Methode")
heading(slide, "Feature Engineering & Modellierung")

# Left column
card_box(slide, Inches(0.5), Inches(2.0), Inches(5.9), Inches(2.5), border_color=CYAN)
add_text(slide, "📐  Feature Engineering",
         Inches(0.7), Inches(2.15), Inches(5.5), Inches(0.4),
         font_size=Pt(14), bold=True, color=CYAN)
fe_items = [
    "Lag-Variablen: 1h, 2h, 3h, 6h, 12h, 24h, 48h, 168h",
    "Gleitende Mittelwerte (rolling mean / std)",
    "Zeitfeatures: Stunde, Wochentag, Monat, Feiertag",
    "Wetter × Preis Interaktionen",
    "Wasserkraft-Indikator (BAFU-Pegel)",
]
bullet_text(slide, fe_items, Inches(0.7), Inches(2.65), Inches(5.5), size=Pt(11.5))

card_box(slide, Inches(0.5), Inches(4.65), Inches(5.9), Inches(1.5), border_color=GREEN)
add_text(slide, "🔀  Train / Val / Test Split",
         Inches(0.7), Inches(4.8), Inches(5.5), Inches(0.4),
         font_size=Pt(14), bold=True, color=GREEN)
add_text(slide,
         "Chronologischer Split – kein Data Leakage. Early Stopping mit Validierungsset. "
         "Modelle täglich neu trainiert und versioniert gespeichert.",
         Inches(0.7), Inches(5.25), Inches(5.5), Inches(0.75),
         font_size=Pt(11.5), color=MUTED)

# Right column
card_box(slide, Inches(6.7), Inches(2.0), Inches(6.0), Inches(1.6), border_color=BLUE)
add_text(slide, "🤖  Model A – Netzlast (kWh)",
         Inches(6.9), Inches(2.15), Inches(5.7), Inches(0.4),
         font_size=Pt(14), bold=True, color=BLUE)
add_text(slide, "Ziel: net_load_kwh  (Bruttolast − PV-Einspeisung)",
         Inches(6.9), Inches(2.6), Inches(5.7), Inches(0.3),
         font_size=Pt(11.5), color=MUTED)
add_text(slide, "XGBRegressor · depth=7 · 300 trees · early stopping · min_child=3",
         Inches(6.9), Inches(2.92), Inches(5.7), Inches(0.5),
         font_size=Pt(11), color=GREEN)

card_box(slide, Inches(6.7), Inches(3.75), Inches(6.0), Inches(1.6), border_color=BLUE)
add_text(slide, "🤖  Model B – EPEX-Preis (€/MWh)",
         Inches(6.9), Inches(3.9), Inches(5.7), Inches(0.4),
         font_size=Pt(14), bold=True, color=BLUE)
add_text(slide, "Ziel: price_eur_mwh  (Day-Ahead EPEX Spotmarkt)",
         Inches(6.9), Inches(4.35), Inches(5.7), Inches(0.3),
         font_size=Pt(11.5), color=MUTED)
add_text(slide, "XGBRegressor · depth=5 · L2=2.0 · min_child=5 · early stopping",
         Inches(6.9), Inches(4.67), Inches(5.7), Inches(0.5),
         font_size=Pt(11), color=GREEN)

card_box(slide, Inches(6.7), Inches(5.5), Inches(6.0), Inches(1.2), border_color=GREEN)
add_text(slide, "💰  Tarif­formel",
         Inches(6.9), Inches(5.65), Inches(5.7), Inches(0.4),
         font_size=Pt(14), bold=True, color=GREEN)
add_text(slide,
         "Model A → Netzpreis  ·  Model B → Energiepreis  →  Gesamttarif in Rp./kWh",
         Inches(6.9), Inches(6.1), Inches(5.7), Inches(0.45),
         font_size=Pt(11.5), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – EVALUATION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
label(slide, "Evaluation")
heading(slide, "Ergebnisse – Modellvergleich")

# -- Model B box (left)
card_box(slide, Inches(0.5), Inches(2.0), Inches(5.9), Inches(4.5), border_color=BLUE)
add_text(slide, "Model B – EPEX Preis (€/MWh)",
         Inches(0.7), Inches(2.15), Inches(5.5), Inches(0.4),
         font_size=Pt(14), bold=True, color=BLUE)
add_text(slide, "Test-Set Evaluation  ·  MAPE",
         Inches(0.7), Inches(2.55), Inches(5.5), Inches(0.3),
         font_size=Pt(10), color=MUTED)

bar_row(slide, "Naive",   1.0,  "23.0 %", Inches(3.05), RED,   left=Inches(0.7), track_w=Inches(3.5))
bar_row(slide, "Linear",  0.40, " 9.2 %", Inches(3.55), YELLOW,left=Inches(0.7), track_w=Inches(3.5))
bar_row(slide, "XGBoost", 0.46, "10.6 %", Inches(4.05), GREEN, left=Inches(0.7), track_w=Inches(3.5))
add_text(slide, "MAPE (niedriger = besser)",
         Inches(0.7), Inches(4.57), Inches(5.5), Inches(0.25),
         font_size=Pt(9), color=MUTED, italic=True)

# Key metrics
add_rect(slide, Inches(0.7), Inches(4.9), Inches(5.5), Pt(1),
         fill_color=BORDER, line_color=None, line_width=Pt(0))
add_text(slide, "XGBoost  MAE = 10.0 €/MWh  ·  RMSE = 20.4 €/MWh",
         Inches(0.7), Inches(5.05), Inches(5.5), Inches(0.4),
         font_size=Pt(12), color=BLUE, bold=True)
add_text(slide, "Verbesserung gegenüber Naive-Baseline:  ×2.2 besser (MAPE)",
         Inches(0.7), Inches(5.5), Inches(5.5), Inches(0.35),
         font_size=Pt(11), color=MUTED)

# -- Model A box (right)
card_box(slide, Inches(6.7), Inches(2.0), Inches(6.0), Inches(4.5), border_color=GREEN)
add_text(slide, "Model A – Netzlast (kWh)",
         Inches(6.9), Inches(2.15), Inches(5.7), Inches(0.4),
         font_size=Pt(14), bold=True, color=GREEN)
add_text(slide, "Test-Set Evaluation  ·  MAPE",
         Inches(6.9), Inches(2.55), Inches(5.7), Inches(0.3),
         font_size=Pt(10), color=MUTED)

bar_row(slide, "Naive",   1.0,  "16.7 %", Inches(3.05), RED,   left=Inches(6.9), track_w=Inches(3.5))
bar_row(slide, "Linear",  0.12, " 2.05 %",Inches(3.55), YELLOW,left=Inches(6.9), track_w=Inches(3.5))
bar_row(slide, "XGBoost", 0.09, " 1.51 %",Inches(4.05), GREEN, left=Inches(6.9), track_w=Inches(3.5))
add_text(slide, "MAPE (niedriger = besser)",
         Inches(6.9), Inches(4.57), Inches(5.7), Inches(0.25),
         font_size=Pt(9), color=MUTED, italic=True)

add_rect(slide, Inches(6.9), Inches(4.9), Inches(5.5), Pt(1),
         fill_color=BORDER, line_color=None, line_width=Pt(0))
add_text(slide, "XGBoost  MAE = 217 kWh  ·  RMSE = 301 kWh",
         Inches(6.9), Inches(5.05), Inches(5.7), Inches(0.4),
         font_size=Pt(12), color=GREEN, bold=True)
add_text(slide, "Ziel MAPE < 8 % erreicht: 1.51 %  ·  ×11 besser als Naive",
         Inches(6.9), Inches(5.5), Inches(5.7), Inches(0.35),
         font_size=Pt(11), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – DEMO / DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
label(slide, "Demo")
heading(slide, "Ergebnis: Dashboard & API")

# Traffic light
card_box(slide, Inches(0.5), Inches(2.0), Inches(5.9), Inches(3.6), border_color=CYAN)
add_text(slide, "🚦  Ampelsystem für Haushalte",
         Inches(0.7), Inches(2.15), Inches(5.5), Inches(0.4),
         font_size=Pt(14), bold=True, color=CYAN)

tl = [
    (GREEN,  "Günstig",  "< 15 Rp./kWh  ·  Jetzt Waschmaschine starten!"),
    (YELLOW, "Mittel",   "15 – 22 Rp./kWh  ·  Normal verbrauchen"),
    (RED,    "Teuer",    "> 22 Rp./kWh  ·  Verbrauch reduzieren"),
]
for i, (col, lbl, val) in enumerate(tl):
    y = Inches(2.75) + i * Inches(0.88)
    add_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.75),
             fill_color=CARD, line_color=col, line_width=Pt(0.75))
    add_rect(slide, Inches(0.85), y + Inches(0.22), Inches(0.32), Inches(0.32),
             fill_color=col, line_color=None, line_width=Pt(0))
    add_text(slide, lbl, Inches(1.3), y + Inches(0.1), Inches(4.7), Inches(0.3),
             font_size=Pt(13), bold=True, color=WHITE)
    add_text(slide, val, Inches(1.3), y + Inches(0.42), Inches(4.7), Inches(0.28),
             font_size=Pt(11), color=MUTED)

# API endpoints
card_box(slide, Inches(6.7), Inches(2.0), Inches(6.0), Inches(2.2), border_color=BLUE)
add_text(slide, "🚀  REST API Endpoints",
         Inches(6.9), Inches(2.15), Inches(5.7), Inches(0.4),
         font_size=Pt(14), bold=True, color=BLUE)
endpoints = [
    ("GET",  "/api/forecast",         GREEN),
    ("GET",  "/api/price-history",    GREEN),
    ("POST", "/api/predict   (JWT)",  BLUE),
    ("POST", "/api/training/trigger", BLUE),
    ("POST", "/api/backfill/trigger", BLUE),
]
for i, (method, ep, col) in enumerate(endpoints):
    y = Inches(2.65) + i * Inches(0.3)
    add_text(slide, method, Inches(6.9), y, Inches(0.7), Inches(0.28),
             font_size=Pt(11), bold=True, color=col)
    add_text(slide, ep, Inches(7.65), y, Inches(4.8), Inches(0.28),
             font_size=Pt(11), color=MUTED)

# Admin dashboard card
card_box(slide, Inches(6.7), Inches(4.35), Inches(6.0), Inches(1.25), border_color=GREEN)
add_text(slide, "📋  Admin Dashboard",
         Inches(6.9), Inches(4.5), Inches(5.7), Inches(0.4),
         font_size=Pt(14), bold=True, color=GREEN)
add_text(slide,
         "Echtzeit DB-Status, Airflow DAG-Übersicht, Modell-Metriken, "
         "Backfill-Trigger und Datentabellen-Explorer.",
         Inches(6.9), Inches(4.95), Inches(5.7), Inches(0.5),
         font_size=Pt(11.5), color=MUTED)

# API response example
add_rect(slide, Inches(0.5), Inches(5.75), Inches(12.2), Inches(0.9),
         fill_color=RGBColor(0x0d, 0x15, 0x26), line_color=BORDER, line_width=Pt(0.75))
add_text(slide,
         '→  GET /api/forecast   |   { "gesamttarif_rp_kwh": 18.4,  "price_level": "medium",  "energiepreis_rp_kwh": 9.2,  "netzpreis_rp_kwh": 9.2 }',
         Inches(0.75), Inches(5.88), Inches(11.8), Inches(0.6),
         font_size=Pt(11), color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – HERAUSFORDERUNGEN
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
label(slide, "Herausforderungen")
heading(slide, "Was uns gefordert hat")

challenges = [
    ("🔌", "API Rate Limits",
     "ENTSO-E erlaubt nur 400 Calls/Tag. Intelligentes Caching, Backfill-Management und Monitoring mit api_call_log Tabelle notwendig.",
     RED),
    ("📉", "Datenqualität",
     "Fehlende Werte, verschiedene Zeitzonen (UTC vs. CET), unterschiedliche Granularität (15 min vs. 1 h) mussten harmonisiert werden.",
     YELLOW),
    ("💡", "Preis-Regime-Shifts",
     "Post-Energiekrise: Preise fielen von 500+ auf <50 €/MWh. Modelle trainiert auf Krisenpreisen generalisieren schlecht.",
     YELLOW),
    ("🔗", "Datenintegration",
     "6 APIs mit unterschiedlichen Formaten, Authentifizierungen und Strukturen in eine konsistente Zeitreihen-DB bringen.",
     BLUE),
    ("🏗️", "Infrastruktur",
     "Docker-Compose-Stack mit Airflow, TimescaleDB, API und Frontend – reproduzierbar, stabil und wartbar zu halten war aufwändig.",
     CYAN),
    ("✅", "Gelöst durch",
     "Robuste Collector-Basisklasse, Retry-Logik, Rate-Limit-Monitoring, chronologischer Train-Split und regularisiertes XGBoost.",
     GREEN),
]

col_w = Inches(4.0)
col_h = Inches(1.45)
gap_x = Inches(0.23)
gap_y = Inches(0.2)
top_c = Inches(2.0)

for i, (icon, title, desc, color) in enumerate(challenges):
    row = i // 3
    col = i  % 3
    lx  = Inches(0.5) + col * (col_w + gap_x)
    ly  = top_c + row * (col_h + gap_y)
    card_box(slide, lx, ly, col_w, col_h, border_color=color)
    add_text(slide, icon + "  " + title,
             lx + Inches(0.2), ly + Inches(0.15), col_w - Inches(0.25), Inches(0.38),
             font_size=Pt(13), bold=True, color=WHITE)
    add_text(slide, desc,
             lx + Inches(0.2), ly + Inches(0.55), col_w - Inches(0.25), Inches(0.78),
             font_size=Pt(10.5), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – LESSONS LEARNED
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
label(slide, "Lessons Learned")
heading(slide, "Was wir gelernt haben")

lessons = [
    ("🏗️", "Pipeline First",
     "Eine robuste, automatisierte Datenpipeline ist die Grundlage – ohne saubere Daten bringt kein Modell etwas."),
    ("📅", "Zeitreihen ≠ Tabellendaten",
     "Data Leakage durch falsches Splitting ist fatal. Chronologischer Split und Walk-Forward-Validation sind zwingend."),
    ("📊", "Baselines nicht unterschätzen",
     "Linear Regression erreichte MAPE 9.2 % beim Preis – nahezu identisch zu XGBoost. Feature Engineering schlägt Komplexität."),
    ("🐳", "Docker von Anfang an",
     "\"Works on my machine\" war kein Problem. Docker Compose machte das Onboarding neuer Teammitglieder trivial."),
    ("🌊", "Domänenwissen zählt",
     "BAFU-Wasserpegel als Feature für Wasserkraft-Einspeisung war eine domänen­spezifische Entscheidung mit messbarem Effekt."),
    ("🤝", "Modulare Architektur",
     "Klare Trennung ETL / DB / Modell / API / Frontend ermöglichte paralleles Arbeiten ohne Konflikte im Team."),
]

col_w = Inches(4.0)
col_h = Inches(1.45)
gap_x = Inches(0.23)
gap_y = Inches(0.2)
top_l = Inches(2.0)

for i, (icon, title, desc) in enumerate(lessons):
    row = i // 3
    col = i  % 3
    lx  = Inches(0.5) + col * (col_w + gap_x)
    ly  = top_l + row * (col_h + gap_y)
    card_box(slide, lx, ly, col_w, col_h, border_color=CYAN if i == 5 else BORDER)
    add_text(slide, icon,  lx + Inches(0.2), ly + Inches(0.15), Inches(0.5), Inches(0.4),
             font_size=Pt(22))
    add_text(slide, title, lx + Inches(0.7), ly + Inches(0.18), col_w - Inches(0.85), Inches(0.4),
             font_size=Pt(13), bold=True, color=WHITE)
    add_text(slide, desc,  lx + Inches(0.2), ly + Inches(0.62), col_w - Inches(0.3), Inches(0.75),
             font_size=Pt(10.5), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – FUTURE WORK & FAZIT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, W, Pt(4), fill_color=BLUE, line_color=None, line_width=Pt(0))
add_rect(slide, 0, H - Pt(4), W, Pt(4), fill_color=CYAN, line_color=None, line_width=Pt(0))

label(slide, "Future Work & Fazit")

# Left – Future Work
add_text(slide, "Future Work", Inches(0.5), Inches(0.85), Inches(5.9), Inches(0.45),
         font_size=Pt(22), bold=True, color=WHITE)

future = [
    "LSTM / Transformer für komplexere zeitliche Abhängigkeiten",
    "Probabilistische Vorhersage mit Konfidenzintervallen",
    "Geräte-Scheduler (Waschmaschine, EV, Wärmepumpe)",
    "Push-Benachrichtigungen bei günstigen Preisfenstern",
    "Ausweitung auf andere Schweizer Städte (Bern, Zürich)",
    "Negative Preise besser modellieren (Übereinspeisung)",
]
bullet_text(slide, future, Inches(0.5), Inches(1.55), Inches(5.9), size=Pt(13))

# Right – Fazit
add_text(slide, "Fazit", Inches(7.0), Inches(0.85), Inches(5.9), Inches(0.45),
         font_size=Pt(22), bold=True, color=WHITE)

# Achievement boxes
achvs = [
    (GREEN, "✅  Vollständige Big-Data-Plattform",
     "Von API-Kollektion über automatisiertes Training bis zum Dashboard – komplett dockerisiert und produktionsreif."),
    (BLUE,  "📉  Starke Modellresultate",
     "MAPE 10.6 % (Preis) · 1.51 % (Last). Model A übertrifft das Qualitätsziel von < 8 % deutlich."),
]
for i, (col, title, desc) in enumerate(achvs):
    y = Inches(1.55) + i * Inches(1.85)
    card_box(slide, Inches(7.0), y, Inches(5.8), Inches(1.65), border_color=col)
    add_text(slide, title, Inches(7.2), y + Inches(0.18), Inches(5.4), Inches(0.42),
             font_size=Pt(14), bold=True, color=col)
    add_text(slide, desc,  Inches(7.2), y + Inches(0.65), Inches(5.4), Inches(0.85),
             font_size=Pt(12), color=MUTED)

# Big closing statement
add_rect(slide, Inches(7.0), Inches(5.35), Inches(5.8), Inches(1.25),
         fill_color=RGBColor(0x0d, 0x19, 0x30), line_color=CYAN, line_width=Pt(1.5))
add_text(slide, "🏆  BigData · SmallPrice",
         Inches(7.2), Inches(5.5), Inches(5.4), Inches(0.5),
         font_size=Pt(20), bold=True, color=CYAN)
add_text(slide, "Günstig Strom verbrauchen – datengetrieben.",
         Inches(7.2), Inches(5.97), Inches(5.4), Inches(0.4),
         font_size=Pt(13), color=MUTED, italic=True)


# ── Save ───────────────────────────────────────────────────────────────────────
out = "docs/presentation/BigDataSmallPrice_Praesentation.pptx"
prs.save(out)
print(f"Saved → {out}")
